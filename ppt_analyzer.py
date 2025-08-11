import argparse
import json
import os
import re
import time
import random
import tempfile
import logging
import traceback
from pptx import Presentation
from PIL import Image
import pytesseract
import google.generativeai as genai
from google.api_core import exceptions

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler()]
)
logger = logging.getLogger('ppt_analyzer')

# Configuration
GENAI_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GENAI_API_KEY:
    logger.error("GOOGLE_API_KEY environment variable not set")
    raise ValueError("GOOGLE_API_KEY environment variable not set")

genai.configure(api_key=GENAI_API_KEY)
MODEL = genai.GenerativeModel('gemini-1.5-flash-latest')

FINAL_PROMPT = """
Analyze this PowerPoint presentation for factual/logical inconsistencies with focus on:
1. Numerical conflicts (values, percentages, timeframes)
2. Contradictory claims (descriptions, comparisons)
3. Timeline mismatches
4. Definitional inconsistencies
5. Internal slide consistency (totals vs breakdowns)
6. Cross-slide comparisons of similar metrics
7. Baseline omissions in comparative claims

Return as valid JSON ONLY with structure:
{
    "inconsistencies": [{
        "type": "CATEGORY",
        "slides": [X, Y, ...],
        "description": "Explanation with context",
        "evidence": ["Full context quote from slide X", "Full context quote from slide Y"],
        "severity": "Critical/High/Medium/Low"
    }]
}

CRITICAL = Major numerical conflicts or contradictory core claims
HIGH = Significant inconsistencies affecting key messages
MEDIUM = Minor numerical mismatches or missing context
LOW = Presentation inconsistencies without material impact

If no issues: {"inconsistencies": []}
"""

def extract_pptx_content(pptx_path):
    """Enhanced extraction with detailed logging"""
    try:
        logger.info(f"Opening presentation: {pptx_path}")
        prs = Presentation(pptx_path)
        slide_data = {}
        logger.info(f"Presentation contains {len(prs.slides)} slides")
        
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            content = []
            logger.debug(f"Processing slide {slide_num}")
            
            # Preserve slide title if exists
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
                content.append(f"TITLE: {title}")
                logger.debug(f" - Title: {title}")
            
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue  # Skip title since we already captured it
                    
                # Text frames
                if shape.has_text_frame:
                    text = " | ".join(p.text for p in shape.text_frame.paragraphs if p.text)
                    if text:
                        content.append(text)
                        logger.debug(f" - Text: {text[:50]}...")
                
                # Tables
                if shape.has_table:
                    table_text = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_text.append(" | ".join(row_data))
                    table_content = "TABLE: " + "\n".join(table_text)
                    content.append(table_content)
                    logger.debug(f" - Table: {table_content[:50]}...")
                
                # Images
                if shape.shape_type == 13:  # Picture
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                        shape.image.save(tmp.name)
                        try:
                            img = Image.open(tmp.name)
                            ocr_text = pytesseract.image_to_string(img).strip()
                            if ocr_text:
                                img_content = "IMAGE: " + ocr_text
                                content.append(img_content)
                                logger.debug(f" - Image text: {ocr_text[:50]}...")
                        except Exception as e:
                            logger.warning(f"OCR error on slide {slide_num}: {str(e)}")
                        finally:
                            os.unlink(tmp.name)
            
            slide_data[slide_num] = "\n".join(content)
            logger.info(f"Extracted {len(content)} elements from slide {slide_num}")
        
        return slide_data
    
    except Exception as e:
        logger.error(f"Extraction failed: {str(e)}")
        logger.error(traceback.format_exc())
        raise

def analyze_slides(slide_data):
    """Enhanced analysis with detailed diagnostics"""
    try:
        # Build full presentation content
        full_content = "\n\n".join(
            [f"--- SLIDE {num} ---\n{content}" 
             for num, content in slide_data.items()]
        )
        
        # Log token estimate
        token_estimate = len(full_content) // 4  # 1 token ≈ 4 characters
        logger.info(f"Sending content to Gemini ({token_estimate} estimated tokens)")
        
        # API call with retries
        for attempt in range(3):
            try:
                logger.info(f"API attempt {attempt+1}/3")
                response = MODEL.generate_content([FINAL_PROMPT, full_content])
                logger.info("Received response from Gemini API")
                
                # Log response for debugging
                with open("gemini_response.txt", "w") as f:
                    f.write(response.text)
                logger.info("Saved raw response to gemini_response.txt")
                
                # Try to extract JSON
                start_idx = response.text.find('{')
                end_idx = response.text.rfind('}')
                if start_idx == -1 or end_idx == -1:
                    logger.error("No JSON found in response")
                    return []
                
                json_str = response.text[start_idx:end_idx+1]
                logger.debug(f"Extracted JSON: {json_str[:200]}...")
                
                try:
                    data = json.loads(json_str)
                    inconsistencies = data.get("inconsistencies", [])
                    logger.info(f"Found {len(inconsistencies)} inconsistencies in response")
                    return inconsistencies
                except json.JSONDecodeError as e:
                    logger.error(f"JSON decode error: {str(e)}")
                    return []
                
            except exceptions.ResourceExhausted as e:
                wait = 10 * (2 ** attempt) + random.uniform(0, 5)
                logger.warning(f"API quota exceeded. Retrying in {wait:.1f}s...")
                time.sleep(wait)
            except Exception as e:
                logger.error(f"API error: {str(e)}")
                logger.error(traceback.format_exc())
                return []
        
        logger.error("API call failed after 3 attempts")
        return []
    
    except Exception as e:
        logger.error(f"Analysis failed: {str(e)}")
        logger.error(traceback.format_exc())
        return []

def main():
    parser = argparse.ArgumentParser(description="PPTX Consistency Analyzer")
    parser.add_argument("input", help="Path to PPTX file")
    parser.add_argument("-o", "--output", help="Output JSON file")
    parser.add_argument("-v", "--verbose", action="store_true", help="Enable debug logging")
    args = parser.parse_args()
    
    if args.verbose:
        logger.setLevel(logging.DEBUG)
    
    logger.info(f"Starting analysis of: {args.input}")
    
    # Extraction Phase
    logger.info("Extracting content from slides...")
    slide_data = extract_pptx_content(args.input)
    logger.info(f"Extracted content from {len(slide_data)} slides")
    
    # Analysis Phase
    logger.info("Analyzing for inconsistencies...")
    inconsistencies = analyze_slides(slide_data)
    logger.info(f"Found {len(inconsistencies)} potential inconsistencies")
    
    # Output Results
    output = {
        "statistics": {
            "total_slides": len(slide_data),
            "issues_found": len(inconsistencies),
            "analysis_time": time.strftime("%Y-%m-%d %H:%M:%S")
        },
        "inconsistencies": inconsistencies
    }
    
    if args.output:
        with open(args.output, "w") as f:
            json.dump(output, f, indent=2)
        logger.info(f"Results saved to {args.output}")
    else:
        print(json.dumps(output, indent=2))
    
    logger.info("Analysis complete")

if __name__ == "__main__":
    main()